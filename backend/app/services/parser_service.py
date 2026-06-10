"""文档解析服务"""
import os
from typing import Optional
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from bs4 import BeautifulSoup
import markdown
import chardet


class ParserService:
    def __init__(self):
        self.supported_types = {
            "pdf": self._parse_pdf,
            "word": self._parse_word,
            "powerpoint": self._parse_powerpoint,
            "text": self._parse_text,
            "markdown": self._parse_markdown,
            "html": self._parse_html,
        }

    def parse(self, file_path: str, file_type: str) -> Optional[str]:
        """解析文档"""
        parser = self.supported_types.get(file_type)
        if not parser:
            return None
        
        try:
            return parser(file_path)
        except Exception as e:
            print(f"解析文档失败: {e}")
            return None

    def _detect_encoding(self, file_path: str) -> str:
        """检测文件编码"""
        with open(file_path, 'rb') as f:
            raw_data = f.read(10000)
            result = chardet.detect(raw_data)
            return result['encoding'] or 'utf-8'

    def _parse_pdf(self, file_path: str) -> str:
        """解析PDF文件"""
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text.strip()

    def _parse_word(self, file_path: str) -> str:
        """解析Word文档"""
        doc = DocxDocument(file_path)
        text = []
        for paragraph in doc.paragraphs:
            text.append(paragraph.text)
        
        # 提取表格内容
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text.append(cell.text)
        
        return "\n".join(text).strip()

    def _parse_powerpoint(self, file_path: str) -> str:
        """解析PowerPoint文档"""
        prs = Presentation(file_path)
        text = []
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        
        return "\n".join(text).strip()

    def _parse_text(self, file_path: str) -> str:
        """解析纯文本文件"""
        encoding = self._detect_encoding(file_path)
        with open(file_path, 'r', encoding=encoding) as f:
            return f.read().strip()

    def _parse_markdown(self, file_path: str) -> str:
        """解析Markdown文件"""
        encoding = self._detect_encoding(file_path)
        with open(file_path, 'r', encoding=encoding) as f:
            md_content = f.read()
        
        # 转换为HTML，然后提取纯文本
        html = markdown.markdown(md_content)
        soup = BeautifulSoup(html, 'html.parser')
        return soup.get_text().strip()

    def _parse_html(self, file_path: str) -> str:
        """解析HTML文件"""
        encoding = self._detect_encoding(file_path)
        with open(file_path, 'r', encoding=encoding) as f:
            html_content = f.read()
        
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # 移除脚本和样式
        for script in soup(["script", "style"]):
            script.decompose()
        
        return soup.get_text().strip()

    def get_document_metadata(self, file_path: str, file_type: str) -> dict:
        """获取文档元数据"""
        metadata = {
            "file_size": os.path.getsize(file_path),
            "created_time": os.path.getctime(file_path),
            "modified_time": os.path.getmtime(file_path),
        }
        
        # 根据文件类型获取特定元数据
        if file_type == "pdf":
            try:
                reader = PdfReader(file_path)
                metadata["page_count"] = len(reader.pages)
                if reader.metadata:
                    metadata["author"] = reader.metadata.author
                    metadata["title"] = reader.metadata.title
            except:
                pass
        
        elif file_type == "word":
            try:
                doc = DocxDocument(file_path)
                metadata["paragraph_count"] = len(doc.paragraphs)
                metadata["table_count"] = len(doc.tables)
                if doc.core_properties:
                    metadata["author"] = doc.core_properties.author
                    metadata["title"] = doc.core_properties.title
            except:
                pass
        
        return metadata